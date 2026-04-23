"""Pluggable document parser protocol and implementations for multiple file formats."""

from dataclasses import dataclass
from pathlib import Path
from typing import Protocol


@dataclass
class TextChunk:
    """A chunk of text extracted from a document with metadata."""

    text: str
    metadata: dict
    page_number: int | None = None
    source_file: str = ""


class DocumentParser(Protocol):
    """Protocol for document parsers."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        """Parse document and return list of text chunks."""
        ...


class PDFParser:
    """Parse PDF files using PyMuPDF (fitz)."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        import fitz

        doc = fitz.open(str(file_path))
        chunks = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                chunks.append(
                    TextChunk(
                        text=text,
                        metadata={"source": file_path.name, "file_type": "pdf"},
                        page_number=page_num + 1,
                        source_file=file_path.name,
                    )
                )
        doc.close()
        return chunks


class DOCXParser:
    """Parse DOCX files using python-docx."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        from docx import Document

        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
        return [
            TextChunk(
                text=full_text,
                metadata={"source": file_path.name, "file_type": "docx"},
                source_file=file_path.name,
            )
        ]


class TXTParser:
    """Parse plain text files."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            text = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        except Exception:
            text = ""
            
        if not text.strip():
            return []
            
        return [
            TextChunk(
                text=text,
                metadata={"source": file_path.name, "file_type": "txt"},
                source_file=file_path.name,
            )
        ]


class PPTXParser:
    """Parse PPTX presentations using python-pptx."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text)
        full_text = "\n\n".join(texts)
        return [
            TextChunk(
                text=full_text,
                metadata={"source": file_path.name, "file_type": "pptx"},
                page_number=1,
                source_file=file_path.name,
            )
        ]


class XLSXParser:
    """Parse XLSX spreadsheets using openpyxl."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    texts.append(row_text)
        full_text = "\n".join(texts)
        return [
            TextChunk(
                text=full_text,
                metadata={"source": file_path.name, "file_type": "xlsx"},
                source_file=file_path.name,
            )
        ]


class MarkdownParser:
    """Parse Markdown files."""

    def parse(self, file_path: Path) -> list[TextChunk]:
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            text = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        except Exception:
            text = ""
            
        if not text.strip():
            return []
            
        return [
            TextChunk(
                text=text,
                metadata={"source": file_path.name, "file_type": "md"},
                source_file=file_path.name,
            )
        ]


class ParserRegistry:
    """Registry that maps file extensions to parser instances."""

    _parsers: dict[str, DocumentParser] = {
        ".pdf": PDFParser(),
        ".docx": DOCXParser(),
        ".txt": TXTParser(),
        ".pptx": PPTXParser(),
        ".xlsx": XLSXParser(),
        ".md": MarkdownParser(),
        ".markdown": MarkdownParser(),
    }

    @classmethod
    def get_parser(cls, file_path: Path) -> DocumentParser:
        """Get the appropriate parser for a file extension."""
        ext = file_path.suffix.lower()
        parser = cls._parsers.get(ext)
        if parser is None:
            raise ValueError(f"Unsupported file type: {ext}")
        return parser

    @classmethod
    def parse_file(cls, file_path: Path) -> list[TextChunk]:
        """Parse a document file and return text chunks."""
        parser = cls.get_parser(file_path)
        return parser.parse(file_path)
